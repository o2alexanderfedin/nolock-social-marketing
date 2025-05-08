#!/usr/bin/env python3
"""
Script to add text boxes to PowerPoint slides
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def add_textboxes_to_slide2():
    """Add text boxes to slide 2"""
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    
    # Load the presentation
    print(f"Loading presentation: {ppt_path}")
    prs = Presentation(ppt_path)
    
    # Get slide 2
    slide = prs.slides[1]  # 0-indexed
    
    # Content to add
    summary = "Content authenticity issues and lack of ownership create fundamental digital trust challenges."
    
    bullet_points = [
        "• Content that can be secretly modified",
        "• No guarantee that what you see is what was produced",
        "• No inherent ownership means no responsibility",
        "• Rampant misinformation without accountability"
    ]
    
    quote = "76% of users struggle to identify authentic content in digital spaces"
    
    # Add summary textbox
    print("Adding summary text box")
    summary_textbox = slide.shapes.add_textbox(
        Inches(1),       # left
        Inches(2.5),     # top
        Inches(8),       # width
        Inches(0.75)     # height
    )
    summary_textbox.text_frame.text = summary
    # Format paragraph
    p = summary_textbox.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.font.italic = True
    p.font.size = Pt(14)
    
    # Add bullet points textbox
    print("Adding bullet points text box")
    bullets_textbox = slide.shapes.add_textbox(
        Inches(1),       # left
        Inches(3.5),     # top
        Inches(4),       # width
        Inches(2)        # height
    )
    
    # Add bullet points
    tf = bullets_textbox.text_frame
    tf.text = ""  # Clear default text
    
    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        p.level = 0
        # p.font.size = Pt(12)  # This won't work directly on p
        
    # Add quote textbox
    print("Adding quote text box")
    quote_textbox = slide.shapes.add_textbox(
        Inches(1),       # left
        Inches(6),       # top
        Inches(8),       # width
        Inches(0.75)     # height
    )
    quote_textbox.text_frame.text = f"\"{quote}\""
    # Format paragraph
    p = quote_textbox.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.italic = True
    p.font.size = Pt(12)
    
    # Save the presentation
    print("Saving updated presentation")
    prs.save(ppt_path)
    print("Slide 2 updated successfully!")
    
    return True

def open_powerpoint():
    """Open PowerPoint and the presentation"""
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    os.system(f"open -a 'Microsoft PowerPoint' '{ppt_path}'")
    print("PowerPoint opened with the presentation")

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description='Add text boxes to PowerPoint slide 2')
    parser.add_argument('--open', action='store_true', help='Open PowerPoint after update')
    
    args = parser.parse_args()
    
    # Update slide 2
    success = add_textboxes_to_slide2()
    
    # Open PowerPoint if requested
    if args.open:
        open_powerpoint()