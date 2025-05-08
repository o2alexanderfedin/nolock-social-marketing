#!/usr/bin/env python3
"""
Script to combine all text into a single text box on slides
"""

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import datetime

def extract_markdown_content(slide_num):
    """Extract content from markdown file for a slide"""
    markdown_dir = "/Users/alexanderfedin/Projects/nolock.social/marketing/pitch-deck-investor-full/slides/"
    slides_export_dir = "/Users/alexanderfedin/Projects/nolock.social/marketing/slides_export/"
    
    formatted_num = f"{slide_num:02d}"
    md_path = os.path.join(markdown_dir, f"slide{formatted_num}.md")
    export_path = os.path.join(slides_export_dir, f"slide{formatted_num}_export.md")
    
    # Check if markdown file exists
    if not os.path.exists(md_path):
        print(f"Warning: Markdown file not found for slide {slide_num}: {md_path}")
        return None
    
    # Read markdown content
    with open(md_path, 'r') as f:
        content = f.read()
    
    # Extract title (first # heading)
    title_match = re.search(r'^# (.+?)$', content, re.MULTILINE)
    title = title_match.group(1) if title_match else ""
    
    # Extract subtitle (first ## heading)
    subtitle_match = re.search(r'^## (.+?)$', content, re.MULTILINE)
    subtitle = subtitle_match.group(1) if subtitle_match else ""
    
    # Extract summary (content right after subtitle - usually in italics)
    summary = ""
    summary_match = re.search(r'^## .+?\n\*(.+?)\*', content, re.DOTALL)
    if summary_match:
        summary = summary_match.group(1).strip()
    
    # Also try to get summary from export file if it exists
    if os.path.exists(export_path):
        with open(export_path, 'r') as f:
            export_content = f.read()
        
        export_summary_match = re.search(r'## Summary\n(.*?)(?=\n\n|\Z)', export_content, re.DOTALL)
        if export_summary_match and not summary:
            summary = export_summary_match.group(1).strip()
    
    # Extract bullet points (starting with - or * in main content area)
    bullet_points = []
    
    # Look for common section headers in investor pitch deck
    section_headers = [
        "## Critical Issues:", "## Key Points:", "## Core Technology:", 
        "## Market Analysis:", "## Two-Pronged Approach:", "## Core Technology:",
        "## Three Critical Market Failures:", "## Building a New Trust Layer:",
        "## Serious Vulnerabilities:", "## Fundamental Design Flaws:",
        "## Key Milestones", "## Use of Funds", "## Seeking", "## Funding Request"
    ]
    
    section_match = None
    for header in section_headers:
        if header in content:
            section_match = re.search(re.escape(header), content)
            if section_match:
                break
    
    if section_match:
        section_start = section_match.start()
        section_text = content[section_start:]
        # Find end of section (next heading or end of file)
        end_match = re.search(r'\n\n>|\n\n---', section_text)
        if end_match:
            section_text = section_text[:end_match.start()]
        
        # Extract all bullet points
        lines = section_text.split('\n')
        for line in lines:
            # Match lines starting with bullet indicators and possibly having nested bullets
            if re.match(r'^\s*[-*•]\s+(.+)$', line):
                bullet_text = re.sub(r'^\s*[-*•]\s+', '', line).strip()
                # Remove markdown links but keep the text
                bullet_text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', bullet_text)
                bullet_points.append({
                    "text": bullet_text,
                    "bold": "**" in line,
                    "level": 0 if not line.startswith("  ") else 1,
                    "italic": "*" in line and not line.startswith("*"),
                    "color": "blue" if "Source:" in line or "source:" in line else None
                })
    
    # Extract quote (if any)
    quote = ""
    quote_match = re.search(r'> *(.*?)(?=\n\n|\Z)', content, re.DOTALL)
    if quote_match:
        quote = quote_match.group(1).strip()
        # Clean up the quote
        quote = quote.replace('*', '').strip()
    
    return {
        'title': title,
        'subtitle': subtitle,
        'summary': summary,
        'bullet_points': bullet_points,
        'quote': quote
    }

def combine_text_on_slide(slide_num):
    """Combine all text into a single textbox on a slide"""
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    
    # Create backup of the slide
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_path = f"/Users/alexanderfedin/Projects/nolock.social/marketing/slide_{slide_num:02d}_backup_combined_{timestamp}.pptx"
    
    # Load the presentation
    print(f"Loading presentation: {ppt_path}")
    prs = Presentation(ppt_path)
    
    # Save backup
    print(f"Saving backup to: {backup_path}")
    prs.save(backup_path)
    
    # Extract content
    content = extract_markdown_content(slide_num)
    if not content:
        print(f"Skipping slide {slide_num} - no content found")
        return False
    
    # Get the slide (0-indexed)
    if slide_num > len(prs.slides):
        print(f"Error: Slide {slide_num} does not exist (total slides: {len(prs.slides)})")
        return False
    
    slide = prs.slides[slide_num - 1]
    
    # Remove any previously added textboxes (shapes after original shapes)
    original_shapes = list(slide.shapes)
    print(f"Found {len(original_shapes)} shapes on slide {slide_num}")
    
    # Keep only first 3 shapes (typically title, subtitle, and image)
    if len(original_shapes) > 3:
        print(f"Removing {len(original_shapes) - 3} previously added shapes")
        for i in range(len(original_shapes) - 1, 2, -1):
            shape = original_shapes[i]
            # We can't directly remove shapes, but we can make them invisible
            shape.left = -10000000
            shape.width = 0
            shape.height = 0
    
    # Update title and subtitle
    if len(original_shapes) >= 1 and content['title']:
        title_shape = original_shapes[0]
        print(f"Updating title to: {content['title']}")
        if hasattr(title_shape, 'text_frame'):
            title_shape.text_frame.text = content['title']
    
    if len(original_shapes) >= 2 and content['subtitle']:
        subtitle_shape = original_shapes[1]
        print(f"Updating subtitle to: {content['subtitle']}")
        if hasattr(subtitle_shape, 'text_frame'):
            subtitle_shape.text_frame.text = content['subtitle']
    
    # Create a single text box for all content
    print("Creating single text box for all content")
    content_textbox = slide.shapes.add_textbox(
        Inches(1),       # left
        Inches(2.5),     # top
        Inches(8),       # width
        Inches(4.5)      # height
    )
    
    tf = content_textbox.text_frame
    tf.word_wrap = True
    
    # Add summary
    if content['summary']:
        p = tf.add_paragraph()
        p.text = content['summary']
        p.alignment = PP_ALIGN.LEFT
        
        # Format summary
        run = p.runs[0]
        run.font.italic = True
        run.font.size = Pt(12)
        
        # Add a blank line
        tf.add_paragraph()
    
    # Add bullet points
    if content['bullet_points']:
        for i, point in enumerate(content['bullet_points']):
            p = tf.add_paragraph()
            
            # Format based on level
            if point.get("level", 0) > 0:
                p.text = f"   • {point['text']}"  # Indent sub-bullets
            else:
                p.text = f"• {point['text']}"
            
            # Format run
            run = p.runs[0]
            run.font.bold = point.get("bold", False)
            run.font.italic = point.get("italic", False)
            run.font.size = Pt(11) if point.get("level", 0) == 0 else Pt(10)
            
            if point.get("color") == "blue":
                run.font.color.rgb = RGBColor(0, 0, 255)
                run.font.size = Pt(9)  # Smaller for source references
    
    # Add quote at the end
    if content['quote']:
        # Add a blank line before quote
        tf.add_paragraph()
        
        p = tf.add_paragraph()
        p.text = f"\"{content['quote']}\""
        p.alignment = PP_ALIGN.CENTER
        
        # Format quote
        run = p.runs[0]
        run.font.italic = True
        run.font.size = Pt(11)
    
    # Save the updated presentation
    print("Saving updated presentation")
    prs.save(ppt_path)
    print(f"Slide {slide_num} updated with combined text!")
    
    return True

def update_all_slides():
    """Update all slides with combined text boxes"""
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    
    # Create backup of the presentation
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_path = f"/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL_backup_all_{timestamp}.pptx"
    
    # Load presentation to get slide count
    prs = Presentation(ppt_path)
    slide_count = len(prs.slides)
    
    # Save backup
    print(f"Creating backup of entire presentation: {backup_path}")
    prs.save(backup_path)
    
    # Update each slide
    updated_count = 0
    for slide_num in range(1, slide_count + 1):
        success = combine_text_on_slide(slide_num)
        if success:
            updated_count += 1
    
    print(f"Updated {updated_count} slides with combined text!")
    return updated_count

def open_powerpoint():
    """Open PowerPoint and the presentation"""
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    os.system(f"open -a 'Microsoft PowerPoint' '{ppt_path}'")
    print("PowerPoint opened with the presentation")

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description='Combine text into a single textbox on slides')
    parser.add_argument('--slide', type=int, help='Update a specific slide')
    parser.add_argument('--all', action='store_true', help='Update all slides')
    parser.add_argument('--open', action='store_true', help='Open PowerPoint after update')
    
    args = parser.parse_args()
    
    if args.slide:
        # Update a specific slide
        combine_text_on_slide(args.slide)
    elif args.all:
        # Update all slides
        update_all_slides()
    else:
        print("Please specify either --slide or --all")
    
    # Open PowerPoint if requested
    if args.open:
        open_powerpoint()