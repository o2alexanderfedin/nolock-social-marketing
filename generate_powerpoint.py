#!/usr/bin/env python3
"""
PowerPoint Generator for NoLock Social Pitch Decks
Creates PowerPoint presentations from markdown slide content
"""

import os
import re
import glob
from datetime import datetime
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Error: python-pptx package is not installed.")
    print("Please install it with: pip install python-pptx")
    exit(1)

# Define paths
MARKETING_DIR = "/Users/alexanderfedin/Projects/nolock.social/marketing"
BASE_PATH = os.path.join(MARKETING_DIR, "pitch-decks/customer-partner")
SIMPLIFIED_DIR = os.path.join(BASE_PATH, "slides-simplified")
DETAILED_DIR = os.path.join(BASE_PATH, "slides")
IMAGES_DIR = os.path.join(BASE_PATH, "images")

SIMPLIFIED_OUTPUT = os.path.join(MARKETING_DIR, "NoLock_Partner_Simplified.pptx")
DETAILED_OUTPUT = os.path.join(MARKETING_DIR, "NoLock_Partner_Detailed.pptx")

def extract_title(content):
    """Extract the slide title from markdown content."""
    match = re.search(r'# Slide \d+: (.*?)(?:\(|$)', content)
    if match:
        return match.group(1).strip()
    return "NoLock Social"

def extract_content(content):
    """Extract the slide content from markdown content."""
    # Extract content between triple backticks
    match = re.search(r'```\s*(.*?)\s*```', content, re.DOTALL)
    if match:
        return match.group(1).strip()
    return ""

def extract_design_elements(content):
    """Extract design elements from markdown content."""
    match = re.search(r'## Design Elements\s*(.*?)(?:##|\Z)', content, re.DOTALL)
    if match:
        return match.group(1).strip()
    return ""

def extract_pitch_notes(content):
    """Extract pitch notes from markdown content."""
    match = re.search(r'## Pitch Notes\s*(.*?)(?:##|<!--|$)', content, re.DOTALL)
    if match:
        return match.group(1).strip()
    return ""

def parse_bullet_points(content):
    """Parse bullet points from content."""
    bullet_points = []
    for line in content.split('\n'):
        line = line.strip()
        if line.startswith('•') or line.startswith('-'):
            # Remove the bullet character and whitespace
            text = line[1:].strip()
            bullet_points.append(text)
        elif line and not line.isspace():
            # Non-empty line without bullet
            bullet_points.append(line)
    return bullet_points

def create_presentation(slides_dir, output_path, presentation_type):
    """Create a PowerPoint presentation from markdown slides."""
    print(f"Creating {presentation_type} presentation...")
    
    # Create a new PowerPoint presentation
    prs = Presentation()
    
    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "NoLock Social"
    subtitle.text = f"{presentation_type} Partner Pitch Deck\nGenerated on {datetime.now().strftime('%Y-%m-%d')}"
    
    # Process all 21 slides
    for slide_num in range(1, 22):
        slide_num_str = f"{slide_num:02d}"
        slide_path = os.path.join(slides_dir, f"slide{slide_num_str}.md")
        
        # Check if slide file exists
        if not os.path.exists(slide_path):
            print(f"Warning: Slide file not found: {slide_path}")
            continue
        
        # Read slide content
        with open(slide_path, 'r') as f:
            slide_content = f.read()
        
        # Extract slide components
        title_text = extract_title(slide_content)
        content_text = extract_content(slide_content)
        design_elements = extract_design_elements(slide_content)
        pitch_notes = extract_pitch_notes(slide_content)
        
        # Determine if content has bullet points
        bullet_points = parse_bullet_points(content_text)
        
        # Choose appropriate slide layout
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title = slide.shapes.title
        title.text = title_text
        
        # Add content
        content = slide.placeholders[1]
        tf = content.text_frame
        
        if bullet_points:
            # Add bullet points
            for i, point in enumerate(bullet_points):
                p = tf.add_paragraph()
                p.text = point
                p.level = 0  # Level 0 is main bullet point
        else:
            # Add regular text
            tf.text = content_text
        
        # Add notes with design elements and pitch notes
        notes = slide.notes_slide.notes_text_frame
        notes.text = f"DESIGN ELEMENTS:\n\n{design_elements}\n\nPITCH NOTES:\n\n{pitch_notes}"
        
        print(f"Added slide {slide_num}: {title_text}")
    
    # Save the presentation
    prs.save(output_path)
    print(f"Saved presentation to: {output_path}")
    return True

def main():
    """Main function to create both presentations."""
    print("=" * 60)
    print("  NoLock Social PowerPoint Generator")
    print("=" * 60)
    
    # Check if python-pptx is installed
    try:
        from pptx import Presentation
    except ImportError:
        print("Error: python-pptx package is not installed.")
        print("Please install it with: pip install python-pptx")
        return False
    
    # Create simplified presentation
    success1 = create_presentation(
        SIMPLIFIED_DIR,
        SIMPLIFIED_OUTPUT,
        "Simplified"
    )
    
    # Create detailed presentation
    success2 = create_presentation(
        DETAILED_DIR,
        DETAILED_OUTPUT,
        "Detailed"
    )
    
    if success1 and success2:
        print("\n" + "=" * 60)
        print("✅ PowerPoint presentations generated successfully!")
        print("=" * 60)
        print("\nOutput files:")
        print(f"- {SIMPLIFIED_OUTPUT}")
        print(f"- {DETAILED_OUTPUT}")
        print("\nFeatures included:")
        print("- Proper slide layouts and formatting")
        print("- Bullet points formatted correctly")
        print("- Design elements included in notes")
        print("- Pitch notes included in presenter view")
        print("- Title slide with generation date")
        return True
    else:
        print("\nError: Failed to generate one or both presentations.")
        return False

if __name__ == "__main__":
    main()