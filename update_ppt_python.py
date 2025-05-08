#!/usr/bin/env python3
"""
Script to update PowerPoint presentation with content from markdown files
"""

import os
import re
import sys
import glob
from pptx import Presentation
import datetime

def extract_title_subtitle(markdown_path):
    """Extract title and subtitle from markdown file"""
    with open(markdown_path, 'r') as f:
        content = f.read()
    
    # Extract title (first # heading)
    title_match = re.search(r'^# (.+?)$', content, re.MULTILINE)
    title = title_match.group(1) if title_match else None
    
    # Extract subtitle (first ## heading)
    subtitle_match = re.search(r'^## (.+?)$', content, re.MULTILINE)
    subtitle = subtitle_match.group(1) if subtitle_match else None
    
    return title, subtitle

def update_slide(presentation, slide_idx, title, subtitle):
    """Update a specific slide with new title and subtitle"""
    if slide_idx >= len(presentation.slides):
        print(f"Error: Slide {slide_idx+1} does not exist in the presentation")
        return False
    
    slide = presentation.slides[slide_idx]
    
    # Find the title and subtitle shapes
    title_shape = None
    subtitle_shape = None
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            # Try to identify title and subtitle shapes
            if title_shape is None:
                title_shape = shape
            elif subtitle_shape is None:
                subtitle_shape = shape
    
    # Update the shapes
    if title_shape and title:
        print(f"  - Updating title: '{title}'")
        title_shape.text = title
    
    if subtitle_shape and subtitle:
        print(f"  - Updating subtitle: '{subtitle}'")
        subtitle_shape.text = subtitle
    
    return True

def main():
    """Main function to update PowerPoint presentation"""
    # File paths
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    markdown_dir = "/Users/alexanderfedin/Projects/nolock.social/marketing/pitch-deck-investor-full/slides/"
    
    # Create backup of original presentation
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_path = f"/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL_backup_{timestamp}.pptx"
    
    # Load the presentation
    print(f"Loading presentation: {ppt_path}")
    prs = Presentation(ppt_path)
    
    # Save backup
    print(f"Saving backup to: {backup_path}")
    prs.save(backup_path)
    
    # Process each slide
    for slide_num in range(1, 22):  # 21 slides
        formatted_num = f"{slide_num:02d}"
        md_path = os.path.join(markdown_dir, f"slide{formatted_num}.md")
        
        if not os.path.exists(md_path):
            print(f"Warning: Markdown file not found for slide {slide_num}: {md_path}")
            continue
        
        print(f"Processing slide {slide_num}...")
        title, subtitle = extract_title_subtitle(md_path)
        
        if title and subtitle:
            update_slide(prs, slide_num-1, title, subtitle)
        else:
            print(f"  - Warning: Could not extract title/subtitle for slide {slide_num}")
    
    # Save the updated presentation
    output_path = ppt_path  # Save to original file
    print(f"Saving updated presentation to: {output_path}")
    prs.save(output_path)
    print("Update complete!")

if __name__ == "__main__":
    main()