#!/usr/bin/env python3
"""
Script to add text boxes to PowerPoint slides with adjusted font sizes
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def add_textboxes_to_slide2():
    """Add text boxes to slide 2 with improved formatting and smaller fonts"""
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    
    # Create backup
    import datetime
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_path = f"/Users/alexanderfedin/Projects/nolock.social/marketing/slide_02_backup_adjusted_{timestamp}.pptx"
    
    # Load the presentation
    print(f"Loading presentation: {ppt_path}")
    prs = Presentation(ppt_path)
    
    # Save backup
    print(f"Saving backup to: {backup_path}")
    prs.save(backup_path)
    
    # Get slide 2
    slide = prs.slides[1]  # 0-indexed
    
    # Remove any previously added textboxes (shapes after the first 3)
    original_shapes = list(slide.shapes)
    if len(original_shapes) > 3:
        print(f"Removing {len(original_shapes) - 3} previously added shapes")
        for i in range(len(original_shapes) - 1, 2, -1):
            shape = original_shapes[i]
            # We can't directly remove shapes, but we can make them invisible
            # by moving them off the slide or setting their size to 0
            shape.left = -10000000
            shape.width = 0
            shape.height = 0
    
    # Content to add
    summary = "Content authenticity issues and lack of ownership create fundamental digital trust challenges."
    
    # Bullet points with formatting information
    bullet_points = [
        {
            "text": "Content that can be secretly modified",
            "bold": True,
            "level": 0
        },
        {
            "text": "No guarantee that what you see is what was produced",
            "bold": False,
            "level": 1
        },
        {
            "text": "Source: Reuters Digital News Report",
            "bold": False,
            "italic": True,
            "level": 1,
            "color": "blue"
        },
        {
            "text": "No inherent ownership means no responsibility",
            "bold": True,
            "level": 0
        },
        {
            "text": "Rampant misinformation without accountability",
            "bold": False,
            "level": 1
        },
        {
            "text": "Source: World Economic Forum",
            "bold": False,
            "italic": True,
            "level": 1,
            "color": "blue"
        }
    ]
    
    quote = "76% of users struggle to identify authentic content in digital spaces"
    
    # Add summary textbox
    print("Adding summary text box")
    summary_textbox = slide.shapes.add_textbox(
        Inches(1),       # left
        Inches(2.5),     # top
        Inches(8),       # width
        Inches(0.75)     # height
    )
    summary_textbox.text_frame.text = summary
    # Format paragraph
    p = summary_textbox.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.italic = True
    run.font.size = Pt(12)  # Smaller font size
    
    # Add bullet points textbox
    print("Adding bullet points text box")
    bullets_textbox = slide.shapes.add_textbox(
        Inches(1),       # left
        Inches(3.25),    # top
        Inches(7),       # width
        Inches(2.5)      # height
    )
    
    # Add bullet points with proper formatting
    tf = bullets_textbox.text_frame
    tf.text = ""  # Clear default text
    
    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point["text"]
        p.level = point.get("level", 0)
        
        # Format run
        run = p.runs[0]
        run.font.bold = point.get("bold", False)
        run.font.italic = point.get("italic", False)
        run.font.size = Pt(11) if point.get("level", 0) == 0 else Pt(10)  # Smaller font sizes
        
        if point.get("color") == "blue":
            run.font.color.rgb = RGBColor(0, 0, 255)
            run.font.size = Pt(9)  # Even smaller for source references
    
    # Add quote textbox
    print("Adding quote text box")
    quote_textbox = slide.shapes.add_textbox(
        Inches(1),       # left
        Inches(6),       # top
        Inches(8),       # width
        Inches(0.75)     # height
    )
    quote_textbox.text_frame.text = f"\"{quote}\""
    # Format paragraph
    p = quote_textbox.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.italic = True
    run.font.size = Pt(11)  # Smaller font size
    
    # Save the updated presentation
    print("Saving updated presentation")
    prs.save(ppt_path)
    print("Slide 2 updated successfully!")
    
    return True

def open_powerpoint():
    """Open PowerPoint and the presentation"""
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    os.system(f"open -a 'Microsoft PowerPoint' '{ppt_path}'")
    print("PowerPoint opened with the presentation")

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description='Add text boxes to PowerPoint slide 2 with adjusted font sizes')
    parser.add_argument('--open', action='store_true', help='Open PowerPoint after update')
    
    args = parser.parse_args()
    
    # Update slide 2
    success = add_textboxes_to_slide2()
    
    # Open PowerPoint if requested
    if args.open:
        open_powerpoint()