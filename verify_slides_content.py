#!/usr/bin/env python3
"""
Script to verify all slide content against markdown files
Ensures all text and hyperlinks from markdown are present in PowerPoint
"""

import os
import re
import sys
from pptx import Presentation
import time
import datetime
import json
import html
from difflib import SequenceMatcher

class SlideVerifier:
    """Class to verify slide content against markdown files"""
    
    def __init__(self):
        """Initialize the slide verifier"""
        self.ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
        self.markdown_dir = "/Users/alexanderfedin/Projects/nolock.social/marketing/pitch-deck-investor-full/slides/"
        self.verification_log = []
        self.issues_found = 0
        
        # Load the presentation
        print(f"Loading presentation: {self.ppt_path}")
        self.prs = Presentation(self.ppt_path)
        
        # Keep track of missing content to fix later
        self.slides_to_fix = set()
        
    def extract_clean_text(self, text):
        """Extract clean text without markdown formatting"""
        # Remove markdown image links
        text = re.sub(r'!\[.*?\]\(.*?\)', '', text)
        
        # Replace markdown links with just the text
        text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)
        
        # Remove navigation links
        text = re.sub(r'\[← Previous\].*?\[Next →\].*?(?=\n|$)', '', text)
        
        # Remove horizontal rules
        text = re.sub(r'---+', '', text)
        
        # Remove bullets and replace with bullet character
        text = re.sub(r'^\s*[-*]\s+', '• ', text, flags=re.MULTILINE)
        
        # Remove extra whitespace and empty lines
        text = re.sub(r'\n\s*\n', '\n', text)
        text = re.sub(r' +', ' ', text)
        
        # Remove heading markers but keep the text
        text = re.sub(r'^#+\s+', '', text, flags=re.MULTILINE)
        
        # Remove emphasis markers but keep the text
        text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
        text = re.sub(r'\*(.*?)\*', r'\1', text)
        
        # Clean up any remaining whitespace
        text = text.strip()
        
        return text
        
    def extract_meaningful_content(self, markdown_path):
        """Extract meaningful content from markdown file"""
        if not os.path.exists(markdown_path):
            print(f"Warning: Markdown file not found: {markdown_path}")
            return {
                'title': '',
                'subtitle': '',
                'summary': '',
                'bullet_points': [],
                'sections': {},
                'hyperlinks': {},
                'cleaned_text': ''
            }
        
        # Read markdown content
        with open(markdown_path, 'r') as f:
            content = f.read()
        
        # Extract title (first # heading)
        title_match = re.search(r'^# (.+?)$', content, re.MULTILINE)
        title = title_match.group(1) if title_match else ""
        
        # Extract subtitle (first ## heading that isn't a section header)
        subtitle = ""
        subtitle_matches = re.finditer(r'^## ([^:]+?)$', content, re.MULTILINE)
        for match in subtitle_matches:
            potential_subtitle = match.group(1)
            if not any(header in potential_subtitle for header in ["Analysis", "Critical", "Key Points", "Technology", "Addressable", "Approach", "Milestones", "Request", "Obtaining"]):
                subtitle = potential_subtitle
                break
        
        # Extract summary (content right after subtitle - usually in italics)
        summary = ""
        summary_match = re.search(r'^## .+?\n\*(.+?)\*', content, re.DOTALL)
        if summary_match:
            summary = summary_match.group(1).strip()
        
        # Extract sections (## headings with colons)
        sections = {}
        section_matches = re.finditer(r'^## (.+?):$', content, re.MULTILINE)
        for match in section_matches:
            section_name = match.group(1)
            section_start = match.end()
            
            # Find the end of the section
            next_section = re.search(r'^##', content[section_start:], re.MULTILINE)
            horizonal_rule = re.search(r'^---', content[section_start:], re.MULTILINE)
            
            if next_section:
                section_end = section_start + next_section.start()
            elif horizonal_rule:
                section_end = section_start + horizonal_rule.start()
            else:
                section_end = len(content)
            
            section_content = content[section_start:section_end].strip()
            sections[section_name] = section_content
        
        # Extract bullet points (- or * at start of line)
        bullet_points = []
        bullet_matches = re.finditer(r'^[ \t]*[-*]\s+(.+?)$', content, re.MULTILINE)
        for match in bullet_matches:
            bullet_points.append(match.group(1).strip())
        
        # Extract hyperlinks
        hyperlinks = {}
        hyperlink_matches = re.finditer(r'\[([^\]]+)\]\(([^)]+)\)', content)
        for match in hyperlink_matches:
            link_text = match.group(1)
            link_url = match.group(2)
            hyperlinks[link_text] = link_url
        
        # Get cleaned text for simple comparison
        cleaned_text = self.extract_clean_text(content)
        
        return {
            'title': title,
            'subtitle': subtitle,
            'summary': summary,
            'bullet_points': bullet_points,
            'sections': sections,
            'hyperlinks': hyperlinks,
            'cleaned_text': cleaned_text
        }
    
    def extract_slide_text(self, slide):
        """Extract all text from a PowerPoint slide"""
        all_text = []
        
        # Extract text from all shapes with text frames
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
                text = shape.text_frame.text.strip()
                if text:
                    all_text.append(text)
        
        # Join all text with line breaks
        return "\n".join(all_text)
    
    def verify_slide(self, slide_num):
        """Verify content on a specific slide"""
        print(f"\n===== Verifying Slide {slide_num} =====")
        self.verification_log.append(f"\n===== Slide {slide_num} =====")
        
        # Get markdown path
        formatted_num = f"{slide_num:02d}"
        md_path = os.path.join(self.markdown_dir, f"slide{formatted_num}.md")
        
        # Check if markdown file exists
        if not os.path.exists(md_path):
            print(f"Warning: Markdown file not found: {md_path}")
            self.verification_log.append(f"❌ Markdown file not found: {md_path}")
            return False
        
        # Extract content from markdown
        md_content = self.extract_meaningful_content(md_path)
        
        # Get the slide (0-indexed)
        if slide_num > len(self.prs.slides):
            print(f"Error: Slide {slide_num} does not exist (total slides: {len(self.prs.slides)})")
            self.verification_log.append(f"❌ Slide {slide_num} does not exist in PowerPoint")
            return False
        
        slide = self.prs.slides[slide_num - 1]
        
        # Extract text from the slide
        slide_text = self.extract_slide_text(slide)
        
        # Clean up slide text for comparison
        cleaned_slide_text = slide_text.replace('\n', ' ').replace('  ', ' ')
        
        # Verify title
        if md_content['title']:
            if md_content['title'] in slide_text:
                print(f"✅ Title found: {md_content['title']}")
                self.verification_log.append(f"✅ Title found: {md_content['title']}")
            else:
                print(f"❌ Title missing: {md_content['title']}")
                self.verification_log.append(f"❌ Title missing: {md_content['title']}")
                self.issues_found += 1
                self.slides_to_fix.add(slide_num)
        
        # Verify subtitle
        if md_content['subtitle']:
            if md_content['subtitle'] in slide_text:
                print(f"✅ Subtitle found: {md_content['subtitle']}")
                self.verification_log.append(f"✅ Subtitle found: {md_content['subtitle']}")
            else:
                print(f"❌ Subtitle missing: {md_content['subtitle']}")
                self.verification_log.append(f"❌ Subtitle missing: {md_content['subtitle']}")
                self.issues_found += 1
                self.slides_to_fix.add(slide_num)
        
        # Verify summary
        if md_content['summary']:
            # Remove asterisks for italics in markdown
            clean_summary = md_content['summary'].replace('*', '')
            if clean_summary in slide_text:
                print(f"✅ Summary found")
                self.verification_log.append(f"✅ Summary found")
            else:
                missing_percentage = self.get_missing_percentage(clean_summary, slide_text)
                print(f"❌ Summary missing ({missing_percentage}% missing)")
                self.verification_log.append(f"❌ Summary missing ({missing_percentage}% missing)")
                self.verification_log.append(f"   Expected: {clean_summary}")
                self.issues_found += 1
                self.slides_to_fix.add(slide_num)
        
        # Verify section headings
        for section_name, section_content in md_content['sections'].items():
            if section_name in slide_text:
                print(f"✅ Section found: {section_name}")
                self.verification_log.append(f"✅ Section found: {section_name}")
            else:
                print(f"❌ Section missing: {section_name}")
                self.verification_log.append(f"❌ Section missing: {section_name}")
                self.issues_found += 1
                self.slides_to_fix.add(slide_num)
        
        # Verify bullet points
        bullet_issues = 0
        for bullet in md_content['bullet_points']:
            # Clean the bullet point text (remove markdown formatting)
            clean_bullet = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', bullet)
            clean_bullet = clean_bullet.replace('**', '').replace('*', '')
            
            if clean_bullet in slide_text or f"• {clean_bullet}" in slide_text:
                print(f"✅ Bullet point found: {clean_bullet[:40]}..." if len(clean_bullet) > 40 else f"✅ Bullet point found: {clean_bullet}")
                self.verification_log.append(f"✅ Bullet point found")
            else:
                # Check if at least 80% of the content is present
                similarity = self.get_text_similarity(clean_bullet, slide_text)
                if similarity >= 0.8:
                    print(f"⚠️ Bullet point partially found ({similarity:.1%}): {clean_bullet[:40]}..." if len(clean_bullet) > 40 else f"⚠️ Bullet point partially found ({similarity:.1%}): {clean_bullet}")
                    self.verification_log.append(f"⚠️ Bullet point partially found ({similarity:.1%})")
                else:
                    print(f"❌ Bullet point missing: {clean_bullet[:40]}..." if len(clean_bullet) > 40 else f"❌ Bullet point missing: {clean_bullet}")
                    self.verification_log.append(f"❌ Bullet point missing: {clean_bullet}")
                    bullet_issues += 1
        
        if bullet_issues > 0:
            self.issues_found += bullet_issues
            self.slides_to_fix.add(slide_num)
        
        # Verify hyperlinks (look for blue, underlined text)
        hyperlink_issues = 0
        for link_text, url in md_content['hyperlinks'].items():
            # Skip navigation links
            if link_text in ["← Previous", "↑ Overview", "Next →"] or url.endswith(".md") or url.startswith("../"):
                continue
                
            # For image links, skip verification
            if url.endswith(".png"):
                continue
                
            # Clean the link text (remove markdown formatting)
            clean_link_text = link_text.replace('**', '').replace('*', '')
            
            # If it's in the slide text, it should be formatted as a hyperlink (blue and underlined)
            if clean_link_text in slide_text or f"Source: {clean_link_text}" in slide_text:
                blue_text_found = False
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'paragraphs'):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                # Check if this run has blue color
                                if (hasattr(run, 'font') and hasattr(run.font, 'color') and 
                                    hasattr(run.font.color, 'rgb') and run.font.color.rgb and 
                                    clean_link_text in run.text):
                                    blue_text_found = True
                                    break
                            if blue_text_found:
                                break
                    if blue_text_found:
                        break
                
                if blue_text_found:
                    print(f"✅ Hyperlink formatting found: {clean_link_text[:40]}..." if len(clean_link_text) > 40 else f"✅ Hyperlink formatting found: {clean_link_text}")
                    self.verification_log.append(f"✅ Hyperlink formatting found")
                else:
                    print(f"❌ Hyperlink not properly formatted: {clean_link_text[:40]}..." if len(clean_link_text) > 40 else f"❌ Hyperlink not properly formatted: {clean_link_text}")
                    self.verification_log.append(f"❌ Hyperlink not properly formatted: {clean_link_text} -> {url}")
                    hyperlink_issues += 1
            else:
                # If the exact link text isn't found, check for approximate matches
                similarity = self.get_text_similarity(clean_link_text, slide_text)
                if similarity >= 0.7:
                    print(f"⚠️ Hyperlink text partially found ({similarity:.1%}): {clean_link_text[:40]}..." if len(clean_link_text) > 40 else f"⚠️ Hyperlink text partially found ({similarity:.1%}): {clean_link_text}")
                    self.verification_log.append(f"⚠️ Hyperlink text partially found ({similarity:.1%})")
                else:
                    print(f"❌ Hyperlink text missing: {clean_link_text[:40]}..." if len(clean_link_text) > 40 else f"❌ Hyperlink text missing: {clean_link_text}")
                    self.verification_log.append(f"❌ Hyperlink text missing: {clean_link_text} -> {url}")
                    hyperlink_issues += 1
        
        if hyperlink_issues > 0:
            self.issues_found += hyperlink_issues
            self.slides_to_fix.add(slide_num)
        
        # Calculate overall content coverage
        md_cleaned = md_content['cleaned_text'].replace('\n', ' ').replace('  ', ' ')
        content_coverage = self.calculate_content_coverage(md_cleaned, cleaned_slide_text)
        print(f"Content coverage: {content_coverage:.1%}")
        self.verification_log.append(f"Content coverage: {content_coverage:.1%}")
        
        # Add overall verdict
        if content_coverage < 0.8:
            print("❌ Significant content missing from slide")
            self.verification_log.append("❌ Significant content missing from slide")
            self.slides_to_fix.add(slide_num)
        elif content_coverage < 0.95:
            print("⚠️ Some content may be missing from slide")
            self.verification_log.append("⚠️ Some content may be missing from slide")
        else:
            print("✅ Most content present on slide")
            self.verification_log.append("✅ Most content present on slide")
        
        return True
    
    def get_missing_percentage(self, expected_text, actual_text):
        """Calculate approximate percentage of missing text"""
        # Break into words and count how many are missing
        expected_words = set(expected_text.split())
        actual_words = set(actual_text.split())
        
        missing_words = expected_words - actual_words
        
        if len(expected_words) == 0:
            return 0
            
        return round(len(missing_words) * 100 / len(expected_words))
    
    def get_text_similarity(self, text1, text2):
        """Get similarity between two text strings"""
        return SequenceMatcher(None, text1, text2).ratio()
    
    def calculate_content_coverage(self, markdown_text, slide_text):
        """Calculate how much of the markdown content is in the slide"""
        # Convert to lowercase for case-insensitive comparison
        markdown_text = markdown_text.lower()
        slide_text = slide_text.lower()
        
        # Break into words for comparison
        markdown_words = set(markdown_text.split())
        slide_words = set(slide_text.split())
        
        # Count words in common
        common_words = markdown_words.intersection(slide_words)
        
        if len(markdown_words) == 0:
            return 1.0
            
        return len(common_words) / len(markdown_words)
    
    def verify_all_slides(self):
        """Verify all slides in the presentation"""
        # Get total number of slides
        total_slides = len(self.prs.slides)
        print(f"Verifying all {total_slides} slides")
        
        # Process each slide
        for slide_num in range(1, total_slides + 1):
            self.verify_slide(slide_num)
        
        # Save verification log
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        log_path = f"/Users/alexanderfedin/Projects/nolock.social/marketing/slide_verification_log_{timestamp}.txt"
        
        with open(log_path, 'w') as f:
            f.write(f"Slide Verification Report\n")
            f.write(f"Generated: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            f.write(f"Total issues found: {self.issues_found}\n")
            f.write(f"Slides with issues: {sorted(list(self.slides_to_fix))}\n\n")
            f.write("\n".join(self.verification_log))
        
        print(f"\n===== Verification Summary =====")
        print(f"Total issues found: {self.issues_found}")
        print(f"Slides with issues: {sorted(list(self.slides_to_fix))}")
        print(f"Verification log saved to: {log_path}")
        
        return self.issues_found
    
    def open_powerpoint(self, slide_num=None):
        """Open PowerPoint with the presentation at a specific slide"""
        ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
        if slide_num:
            # Use AppleScript to open PowerPoint and go to specific slide
            script = f'''
            tell application "Microsoft PowerPoint"
                activate
                open "{ppt_path}"
                tell active presentation
                    go to slide {slide_num}
                end tell
            end tell
            '''
            os.system(f"osascript -e '{script}'")
            print(f"PowerPoint opened to slide {slide_num}")
        else:
            os.system(f"open -a 'Microsoft PowerPoint' '{ppt_path}'")
            print("PowerPoint opened with the presentation")

def main():
    """Main function"""
    import argparse
    
    parser = argparse.ArgumentParser(description='Verify PowerPoint slides against markdown content')
    parser.add_argument('--slide', type=int, help='Verify a specific slide')
    parser.add_argument('--all', action='store_true', help='Verify all slides')
    parser.add_argument('--open', action='store_true', help='Open PowerPoint after verification')
    
    args = parser.parse_args()
    
    # Create slide verifier
    verifier = SlideVerifier()
    
    # Verify slides
    if args.slide:
        # Verify a specific slide
        verifier.verify_slide(args.slide)
        
        # Open PowerPoint if requested
        if args.open:
            verifier.open_powerpoint(args.slide)
    elif args.all:
        # Verify all slides
        verifier.verify_all_slides()
        
        # Open PowerPoint if requested
        if args.open:
            if verifier.slides_to_fix:
                # Open to the first slide with issues
                verifier.open_powerpoint(min(verifier.slides_to_fix))
            else:
                verifier.open_powerpoint()
    else:
        print("Please specify either --slide or --all")

if __name__ == "__main__":
    main()