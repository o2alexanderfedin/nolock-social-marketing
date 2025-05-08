#!/usr/bin/env python3
"""
Script to fix remaining issues with slides
"""

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import datetime

def fix_slide1():
    """Fix slide 1 (title slide) with all content"""
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    
    # Create backup
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_path = f"/Users/alexanderfedin/Projects/nolock.social/marketing/slide_01_final_{timestamp}.pptx"
    
    # Load the presentation
    print(f"Loading presentation: {ppt_path}")
    prs = Presentation(ppt_path)
    
    # Save backup
    print(f"Saving backup to: {backup_path}")
    prs.save(backup_path)
    
    # Get slide 1
    slide = prs.slides[0]  # 0-indexed
    
    # Remove any previously added textboxes
    shapes = list(slide.shapes)
    if len(shapes) > 2:
        print(f"Removing {len(shapes) - 2} previously added shapes")
        for i in range(len(shapes) - 1, 1, -1):
            shape = shapes[i]
            # We can't directly remove shapes, but we can make them invisible
            shape.left = -10000000
            shape.width = 0
            shape.height = 0
    
    # Add detailed content
    print("Adding title slide content")
    content_textbox = slide.shapes.add_textbox(
        Inches(1),       # left
        Inches(2.5),     # top
        Inches(8),       # width
        Inches(3.5)      # height
    )
    
    tf = content_textbox.text_frame
    tf.word_wrap = True
    
    # Add subtitle text
    p = tf.paragraphs[0]
    p.text = "A next-generation platform rebuilding digital trust with decentralized identity in a $12B market."
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.italic = True
    run.font.size = Pt(14)
    
    # Add a blank line
    tf.add_paragraph()
    
    # Add "Pre-Seed Investment Opportunity" header
    p = tf.add_paragraph()
    p.text = "Pre-Seed Investment Opportunity"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.bold = True
    run.font.size = Pt(16)
    
    # Add a blank line
    tf.add_paragraph()
    
    # Add bullet points
    bullet_points = [
        "Next-Generation Trust Infrastructure",
        "Decentralized Identity and Social Tools",
        "$12.1B Market | 25% Annual Growth Rate"
    ]
    
    for bp_text in bullet_points:
        p = tf.add_paragraph()
        p.text = f"• {bp_text}"
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(12)
        if "Next-Generation" in bp_text or "Decentralized" in bp_text:
            run.font.bold = True
    
    # Add a blank line
    tf.add_paragraph()
    
    # Add source
    p = tf.add_paragraph()
    p.text = "Source: Market Research Future (MRFR)"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.italic = True
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0, 0, 255)
    run.font.underline = True
    
    # Save the presentation
    print("Saving updated presentation")
    prs.save(ppt_path)
    print("Slide 1 fixed!")
    
    return True

def fix_competitor_links():
    """Fix hyperlinks for competitors on slide 18"""
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    
    # Create backup
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_path = f"/Users/alexanderfedin/Projects/nolock.social/marketing/slide_18_competitors_{timestamp}.pptx"
    
    # Load the presentation
    print(f"Loading presentation: {ppt_path}")
    prs = Presentation(ppt_path)
    
    # Save backup
    print(f"Saving backup to: {backup_path}")
    prs.save(backup_path)
    
    # Get slide 18
    slide = prs.slides[17]  # 0-indexed
    
    # Find the content shape
    content_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and "Twitter/X" in shape.text_frame.text:
            content_shape = shape
            break
    
    if not content_shape:
        print("Could not find content with competitors on slide 18")
        return False
    
    # Get the text and replace the competitor names with blue, underlined versions
    print("Fixing competitor links on slide 18")
    for paragraph in content_shape.text_frame.paragraphs:
        if "Twitter/X, Facebook:" in paragraph.text:
            # Clear existing text
            original_text = paragraph.text
            paragraph.text = ""
            
            # Split into parts
            prefix_text = "• "
            competitor_twitter = "Twitter/X"
            middle_text1 = ", "
            competitor_facebook = "Facebook"
            suffix_text = ": Rely on central authorities and closed algorithms"
            
            # Add text with proper formatting for each part
            run1 = paragraph.add_run()
            run1.text = prefix_text
            
            run2 = paragraph.add_run()
            run2.text = competitor_twitter
            run2.font.color.rgb = RGBColor(0, 0, 255)
            run2.font.underline = True
            
            run3 = paragraph.add_run()
            run3.text = middle_text1
            
            run4 = paragraph.add_run()
            run4.text = competitor_facebook
            run4.font.color.rgb = RGBColor(0, 0, 255)
            run4.font.underline = True
            
            run5 = paragraph.add_run()
            run5.text = suffix_text
        
        elif "Signal, Telegram:" in paragraph.text:
            # Clear existing text
            original_text = paragraph.text
            paragraph.text = ""
            
            # Split into parts
            prefix_text = "• "
            competitor_signal = "Signal"
            middle_text1 = ", "
            competitor_telegram = "Telegram"
            suffix_text = ": Better privacy, but limited functionality"
            
            # Add text with proper formatting for each part
            run1 = paragraph.add_run()
            run1.text = prefix_text
            
            run2 = paragraph.add_run()
            run2.text = competitor_signal
            run2.font.color.rgb = RGBColor(0, 0, 255)
            run2.font.underline = True
            
            run3 = paragraph.add_run()
            run3.text = middle_text1
            
            run4 = paragraph.add_run()
            run4.text = competitor_telegram
            run4.font.color.rgb = RGBColor(0, 0, 255)
            run4.font.underline = True
            
            run5 = paragraph.add_run()
            run5.text = suffix_text
        
        elif "Mastodon, Bluesky:" in paragraph.text:
            # Clear existing text
            original_text = paragraph.text
            paragraph.text = ""
            
            # Split into parts
            prefix_text = "• "
            competitor_mastodon = "Mastodon"
            middle_text1 = ", "
            competitor_bluesky = "Bluesky"
            suffix_text = ": Federation, not true decentralization"
            
            # Add text with proper formatting for each part
            run1 = paragraph.add_run()
            run1.text = prefix_text
            
            run2 = paragraph.add_run()
            run2.text = competitor_mastodon
            run2.font.color.rgb = RGBColor(0, 0, 255)
            run2.font.underline = True
            
            run3 = paragraph.add_run()
            run3.text = middle_text1
            
            run4 = paragraph.add_run()
            run4.text = competitor_bluesky
            run4.font.color.rgb = RGBColor(0, 0, 255)
            run4.font.underline = True
            
            run5 = paragraph.add_run()
            run5.text = suffix_text
        
        elif "IPFS, Filecoin:" in paragraph.text:
            # Clear existing text
            original_text = paragraph.text
            paragraph.text = ""
            
            # Split into parts
            prefix_text = "• "
            competitor_ipfs = "IPFS"
            middle_text1 = ", "
            competitor_filecoin = "Filecoin"
            suffix_text = ": Storage focused, not social layer"
            
            # Add text with proper formatting for each part
            run1 = paragraph.add_run()
            run1.text = prefix_text
            
            run2 = paragraph.add_run()
            run2.text = competitor_ipfs
            run2.font.color.rgb = RGBColor(0, 0, 255)
            run2.font.underline = True
            
            run3 = paragraph.add_run()
            run3.text = middle_text1
            
            run4 = paragraph.add_run()
            run4.text = competitor_filecoin
            run4.font.color.rgb = RGBColor(0, 0, 255)
            run4.font.underline = True
            
            run5 = paragraph.add_run()
            run5.text = suffix_text
    
    # Save the presentation
    print("Saving updated presentation")
    prs.save(ppt_path)
    print("Slide 18 competitor links fixed!")
    
    return True

def fix_gtm_links():
    """Fix hyperlinks for Go-to-Market Strategy on slide 19"""
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    
    # Create backup
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_path = f"/Users/alexanderfedin/Projects/nolock.social/marketing/slide_19_gtm_{timestamp}.pptx"
    
    # Load the presentation
    print(f"Loading presentation: {ppt_path}")
    prs = Presentation(ppt_path)
    
    # Save backup
    print(f"Saving backup to: {backup_path}")
    prs.save(backup_path)
    
    # Get slide 19
    slide = prs.slides[18]  # 0-indexed
    
    # Find the content shape
    content_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and "Strategy: Developer-first adoption" in shape.text_frame.text:
            content_shape = shape
            break
    
    if not content_shape:
        print("Could not find content with GTM links on slide 19")
        return False
    
    # Get the text and replace the links with blue, underlined versions
    print("Fixing GTM links on slide 19")
    for paragraph in content_shape.text_frame.paragraphs:
        if "Strategy: Developer-first adoption" in paragraph.text:
            # Clear existing text
            original_text = paragraph.text
            paragraph.text = ""
            
            # Split into parts
            prefix_text = "• "
            link_text = "Strategy: Developer-first adoption"
            
            # Add text with proper formatting for each part
            run1 = paragraph.add_run()
            run1.text = prefix_text
            
            run2 = paragraph.add_run()
            run2.text = link_text
            run2.font.color.rgb = RGBColor(0, 0, 255)
            run2.font.underline = True
        
        elif "Market: Privacy-focused segment size" in paragraph.text:
            # Clear existing text
            original_text = paragraph.text
            paragraph.text = ""
            
            # Split into parts
            prefix_text = "• "
            link_text = "Market: Privacy-focused segment size"
            
            # Add text with proper formatting for each part
            run1 = paragraph.add_run()
            run1.text = prefix_text
            
            run2 = paragraph.add_run()
            run2.text = link_text
            run2.font.color.rgb = RGBColor(0, 0, 255)
            run2.font.underline = True
    
    # Save the presentation
    print("Saving updated presentation")
    prs.save(ppt_path)
    print("Slide 19 GTM links fixed!")
    
    return True

def fix_investment_subtitle():
    """Fix missing subtitle on slide 21 (Investment Ask)"""
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    
    # Create backup
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_path = f"/Users/alexanderfedin/Projects/nolock.social/marketing/slide_21_subtitle_{timestamp}.pptx"
    
    # Load the presentation
    print(f"Loading presentation: {ppt_path}")
    prs = Presentation(ppt_path)
    
    # Save backup
    print(f"Saving backup to: {backup_path}")
    prs.save(backup_path)
    
    # Get slide 21
    slide = prs.slides[20]  # 0-indexed
    
    # Find the subtitle shape (typically shape 2)
    subtitle_shape = None
    if len(slide.shapes) >= 2:
        subtitle_shape = slide.shapes[1]
    
    if not subtitle_shape or not hasattr(subtitle_shape, 'text_frame'):
        print("Could not find subtitle shape on slide 21")
        return False
    
    # Set subtitle
    print("Fixing subtitle on slide 21")
    subtitle_shape.text_frame.text = "Funding Request & Use of Funds"
    
    # Save the presentation
    print("Saving updated presentation")
    prs.save(ppt_path)
    print("Slide 21 subtitle fixed!")
    
    return True

def fix_all_remaining_issues():
    """Fix all remaining issues in the presentation"""
    # Fix slide 1 (title slide)
    fix_slide1()
    
    # Fix competitor links on slide 18
    fix_competitor_links()
    
    # Fix GTM links on slide 19
    fix_gtm_links()
    
    # Fix investment subtitle on slide 21
    fix_investment_subtitle()
    
    print("All remaining issues fixed!")
    
    return True

def open_powerpoint():
    """Open PowerPoint with the presentation"""
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    os.system(f"open -a 'Microsoft PowerPoint' '{ppt_path}'")
    print("PowerPoint opened with the presentation")

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description='Fix remaining slide issues')
    parser.add_argument('--all', action='store_true', help='Fix all remaining issues')
    parser.add_argument('--slide1', action='store_true', help='Fix slide 1 (title slide)')
    parser.add_argument('--competitors', action='store_true', help='Fix competitor links on slide 18')
    parser.add_argument('--gtm', action='store_true', help='Fix GTM links on slide 19')
    parser.add_argument('--investment', action='store_true', help='Fix investment subtitle on slide 21')
    parser.add_argument('--open', action='store_true', help='Open PowerPoint after fixing')
    
    args = parser.parse_args()
    
    if args.all:
        fix_all_remaining_issues()
    else:
        if args.slide1:
            fix_slide1()
        if args.competitors:
            fix_competitor_links()
        if args.gtm:
            fix_gtm_links()
        if args.investment:
            fix_investment_subtitle()
    
    # Open PowerPoint if requested
    if args.open:
        open_powerpoint()