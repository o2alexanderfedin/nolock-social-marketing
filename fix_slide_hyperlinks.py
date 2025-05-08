#!/usr/bin/env python3
"""
Script to fix hyperlinks on slides and verify content
"""

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import datetime

def extract_hyperlinks_from_markdown(slide_num):
    """Extract hyperlinks from markdown file for a specific slide"""
    markdown_dir = "/Users/alexanderfedin/Projects/nolock.social/marketing/pitch-deck-investor-full/slides/"
    
    formatted_num = f"{slide_num:02d}"
    md_path = os.path.join(markdown_dir, f"slide{formatted_num}.md")
    
    if not os.path.exists(md_path):
        print(f"Warning: Markdown file not found for slide {slide_num}: {md_path}")
        return {}
    
    # Read markdown content
    with open(md_path, 'r') as f:
        content = f.read()
    
    # Extract all hyperlinks with their text using regex
    hyperlink_pattern = r'\[([^\]]+)\]\(([^)]+)\)'
    hyperlinks = {}
    
    for match in re.finditer(hyperlink_pattern, content):
        link_text = match.group(1)
        link_url = match.group(2)
        hyperlinks[link_text] = link_url
    
    # Also search for lines containing link text for reference
    lines_with_links = []
    for line in content.split('\n'):
        if re.search(hyperlink_pattern, line):
            cleaned_line = re.sub(r'^\s*[-*•]\s+', '', line).strip()
            lines_with_links.append(cleaned_line)
    
    return {
        'hyperlinks': hyperlinks,
        'lines_with_links': lines_with_links
    }

def verify_slide_content(slide_num):
    """Verify content on a slide and fix hyperlinks"""
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    
    # Create backup
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_path = f"/Users/alexanderfedin/Projects/nolock.social/marketing/slide_{slide_num:02d}_verified_{timestamp}.pptx"
    
    # Load the presentation
    print(f"Loading presentation: {ppt_path}")
    prs = Presentation(ppt_path)
    
    # Save backup
    print(f"Saving backup to: {backup_path}")
    prs.save(backup_path)
    
    # Get hyperlinks from markdown
    hyperlink_data = extract_hyperlinks_from_markdown(slide_num)
    hyperlinks = hyperlink_data.get('hyperlinks', {})
    lines_with_links = hyperlink_data.get('lines_with_links', [])
    
    print(f"Found {len(hyperlinks)} hyperlinks in markdown for slide {slide_num}")
    for text, url in hyperlinks.items():
        print(f"  {text} -> {url}")
    
    # Get the slide (0-indexed)
    if slide_num > len(prs.slides):
        print(f"Error: Slide {slide_num} does not exist (total slides: {len(prs.slides)})")
        return False
    
    slide = prs.slides[slide_num - 1]
    
    # Find all text shapes on the slide
    text_shapes = []
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text_shapes.append(shape)
    
    print(f"Found {len(text_shapes)} text shapes on slide {slide_num}")
    
    # Examine and mark hyperlinks
    links_found = False
    
    for shape in text_shapes:
        if not hasattr(shape, 'text_frame'):
            continue
            
        text_frame = shape.text_frame
        
        # Check each paragraph in the text frame
        for paragraph in text_frame.paragraphs:
            paragraph_text = paragraph.text.strip()
            
            # Check if this paragraph text needs a hyperlink
            potential_link = False
            for link_text in hyperlinks.keys():
                if link_text in paragraph_text:
                    potential_link = True
                    break
            
            # Mark lines that should have hyperlinks
            if potential_link or any(line in paragraph_text for line in lines_with_links):
                # Highlight the text in blue to indicate it should be a hyperlink
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(0, 0, 255)
                    run.font.underline = True
                links_found = True
                print(f"  Marked hyperlink in text: {paragraph_text[:50]}..." if len(paragraph_text) > 50 else f"  Marked hyperlink in text: {paragraph_text}")
    
    # Save the updated presentation
    if links_found:
        print("Saving presentation with marked hyperlinks")
        prs.save(ppt_path)
        print(f"Slide {slide_num} updated with marked hyperlinks!")
    else:
        print(f"No hyperlinks to mark on slide {slide_num}")
    
    return links_found

def check_all_slides():
    """Check all slides for hyperlinks"""
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    
    # Load presentation to get slide count
    prs = Presentation(ppt_path)
    slide_count = len(prs.slides)
    
    # Check each slide
    slides_with_links = []
    for slide_num in range(1, slide_count + 1):
        print(f"\n===== Checking Slide {slide_num} =====")
        links_found = verify_slide_content(slide_num)
        if links_found:
            slides_with_links.append(slide_num)
    
    # Summary
    print("\n===== Summary =====")
    print(f"Slides with hyperlinks: {slides_with_links}")
    
    return slides_with_links

def open_powerpoint():
    """Open PowerPoint and the presentation"""
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    os.system(f"open -a 'Microsoft PowerPoint' '{ppt_path}'")
    print("PowerPoint opened with the presentation")

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description='Verify slide content and mark hyperlinks')
    parser.add_argument('--slide', type=int, help='Check a specific slide')
    parser.add_argument('--all', action='store_true', help='Check all slides')
    parser.add_argument('--open', action='store_true', help='Open PowerPoint after update')
    
    args = parser.parse_args()
    
    if args.slide:
        # Check a specific slide
        verify_slide_content(args.slide)
    elif args.all:
        # Check all slides
        check_all_slides()
    else:
        print("Please specify either --slide or --all")
    
    # Open PowerPoint if requested
    if args.open:
        open_powerpoint()