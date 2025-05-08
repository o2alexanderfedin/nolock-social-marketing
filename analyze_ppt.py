#!/usr/bin/env python3
"""
Script to analyze PowerPoint structure in detail
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def analyze_slide(slide_num):
    """Analyze a specific slide in detail"""
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    
    # Load the presentation
    print(f"Loading presentation: {ppt_path}")
    prs = Presentation(ppt_path)
    
    # Validate slide number
    if slide_num > len(prs.slides):
        print(f"Error: Slide {slide_num} does not exist (total slides: {len(prs.slides)})")
        return
    
    slide = prs.slides[slide_num - 1]
    shapes = list(slide.shapes)
    
    print(f"\n===== DETAILED ANALYSIS OF SLIDE {slide_num} =====")
    print(f"Total shapes: {len(shapes)}")
    
    for i, shape in enumerate(shapes):
        print(f"\nShape {i+1}:")
        print(f"  Name: {shape.name if hasattr(shape, 'name') else 'No name'}")
        print(f"  ID: {shape.shape_id if hasattr(shape, 'shape_id') else 'No ID'}")
        print(f"  Type: {shape.shape_type}")
        
        # Type-specific information
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            ph = shape.placeholder_format
            print(f"  Placeholder Type: {ph.type if hasattr(ph, 'type') else 'Unknown'}")
            print(f"  Placeholder Index: {ph.idx if hasattr(ph, 'idx') else 'Unknown'}")
        
        # Check for text frame
        if hasattr(shape, 'text_frame'):
            print("  Has text frame: Yes")
            if hasattr(shape.text_frame, 'text'):
                text = shape.text_frame.text
                print(f"  Text: \"{text[:100]}{'...' if len(text) > 100 else ''}\"")
                
                # Check paragraphs
                print(f"  Paragraphs: {len(shape.text_frame.paragraphs)}")
                for j, para in enumerate(shape.text_frame.paragraphs):
                    print(f"    Para {j+1}: \"{para.text[:50]}{'...' if len(para.text) > 50 else ''}\"")
            else:
                print("  Text frame has no text attribute")
        else:
            print("  Has text frame: No")
        
        # Position and size information
        if hasattr(shape, 'left') and hasattr(shape, 'top'):
            print(f"  Position: left={shape.left}, top={shape.top}")
        if hasattr(shape, 'width') and hasattr(shape, 'height'):
            print(f"  Size: width={shape.width}, height={shape.height}")

def list_all_slides():
    """List all slides with basic information"""
    ppt_path = "/Users/alexanderfedin/Projects/nolock.social/marketing/REBUILD TRUST IN THE DIGITAL SPACEN LOCK•SOCIAL.pptx"
    
    # Load the presentation
    print(f"Loading presentation: {ppt_path}")
    prs = Presentation(ppt_path)
    
    print(f"\n===== PRESENTATION OVERVIEW =====")
    print(f"Total slides: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        shapes = list(slide.shapes)
        text_shapes = [s for s in shapes if hasattr(s, 'text_frame') and hasattr(s.text_frame, 'text')]
        
        title = "No title"
        if text_shapes and hasattr(text_shapes[0].text_frame, 'text'):
            title = text_shapes[0].text_frame.text
        
        print(f"Slide {i+1}: Title=\"{title[:40]}{'...' if len(title) > 40 else ''}\", {len(shapes)} shapes, {len(text_shapes)} text shapes")

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description='Analyze PowerPoint structure')
    parser.add_argument('--slide', type=int, help='Analyze a specific slide')
    parser.add_argument('--all', action='store_true', help='List all slides')
    
    args = parser.parse_args()
    
    if args.slide:
        analyze_slide(args.slide)
    elif args.all:
        list_all_slides()
    else:
        print("Please specify either --slide or --all")